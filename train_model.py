import os
import pickle
import json
import pdfplumber
from docx import Document
from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.linear_model import LogisticRegression
from sklearn.model_selection import cross_val_score
import mysql.connector
import logging
from dotenv import load_dotenv
import shutil

# Load environment variables
load_dotenv()

# Set up logging
logging.basicConfig(
    level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s"
)

DATASET_DIR = "C:\\Doc"
TEMP_DIR = "TempTraining"
DB_CONFIG = {
    "host": os.getenv("DB_HOST"),
    "user": os.getenv("DB_USER"),
    "password": os.getenv("DB_PASSWORD"),
    "database": os.getenv("DB_NAME"),
}

texts = []
labels = []


def extract_content_from_pdf(file_path):
    text = ""
    tables = []
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text() or ""
                text += page_text + "\n"
                page_tables = page.extract_tables()
                for table in page_tables:
                    cleaned_table = [
                        [cell or "" for cell in row]
                        for row in table
                        if any(cell for cell in row)
                    ]
                    if cleaned_table:
                        tables.append(
                            {"page": page.page_number, "table": cleaned_table}
                        )
    except Exception as e:
        logging.error(f"Failed to read PDF {file_path}: {e}")
    return text.strip(), tables


def extract_text_from_docx(file_path):
    text = ""
    try:
        doc = Document(file_path)
        text = "\n".join(para.text for para in doc.paragraphs if para.text.strip())
    except Exception as e:
        logging.error(f"Failed to read DOCX {file_path}: {e}")
    return text.strip(), []


def extract_text_from_pptx(file_path):
    text = []
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
    except Exception as e:
        logging.error(f"Failed to read PPTX {file_path}: {e}")
    return "\n".join(text).strip(), []


def extract_content(file_path):
    if file_path.lower().endswith(".pdf"):
        return extract_content_from_pdf(file_path)
    elif file_path.lower().endswith(".docx"):
        return extract_text_from_docx(file_path)
    elif file_path.lower().endswith(".pptx"):
        return extract_text_from_pptx(file_path)
    else:
        logging.warning(f"Unsupported file type: {file_path}")
        return "", []


# Create temporary training directory
if os.path.exists(TEMP_DIR):
    shutil.rmtree(TEMP_DIR)
os.makedirs(TEMP_DIR)

# Connect to MySQL
try:
    conn = mysql.connector.connect(**DB_CONFIG)
    cursor = conn.cursor()
except Exception as e:
    logging.error(f"Failed to connect to MySQL: {e}")
    exit(1)

# Clear existing pretrained data
cursor.execute("DELETE FROM files WHERE user_id IS NULL")
conn.commit()

# Walk through dataset
for folder_name in os.listdir(DATASET_DIR):
    folder_path = os.path.join(DATASET_DIR, folder_name)
    if os.path.isdir(folder_path):
        temp_folder = os.path.join(TEMP_DIR, folder_name)
        os.makedirs(temp_folder, exist_ok=True)
        for file_name in os.listdir(folder_path):
            if file_name.lower().endswith((".pdf", ".docx", ".pptx")):
                original_path = os.path.join(folder_path, file_name)
                temp_path = os.path.join(temp_folder, file_name)
                shutil.copy2(original_path, temp_path)
                text, tables = extract_content(temp_path)
                if text or tables:
                    texts.append(text)
                    labels.append(folder_name)
                    cursor.execute(
                        "INSERT INTO files (user_id, filename, folder, text, tables, file_path) VALUES (%s, %s, %s, %s, %s, %s)",
                        (
                            None,
                            file_name,
                            folder_name,
                            text,
                            json.dumps(tables),
                            temp_path,
                        ),
                    )
                    conn.commit()
                    logging.info(f"✓ Processed: {file_name} → {folder_name}")
                else:
                    logging.warning(f"Skipped (no content): {file_name}")

cursor.close()
conn.close()

# Clean up temporary directory
shutil.rmtree(TEMP_DIR)

logging.info(f"Collected {len(texts)} documents from {len(set(labels))} categories.")

if not texts:
    logging.error("No documents processed. Cannot train model.")
    exit(1)

# Train model
try:
    vectorizer = TfidfVectorizer(max_features=5000, stop_words="english")
    X = vectorizer.fit_transform(texts)
    model = LogisticRegression(max_iter=1000, C=1.0)
    model.fit(X, labels)
    scores = cross_val_score(model, X, labels, cv=5, scoring="accuracy")
    logging.info(
        f"Cross-validation accuracy: {scores.mean():.2f} (+/- {scores.std() * 2:.2f})"
    )
    with open("model.pkl", "wb") as f:
        pickle.dump(model, f)
    with open("vectorizer.pkl", "wb") as f:
        pickle.dump(vectorizer, f)
    logging.info("Model training complete. Files saved: model.pkl, vectorizer.pkl")
except Exception as e:
    logging.error(f"Failed to train model: {e}")
    exit(1)
