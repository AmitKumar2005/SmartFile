import os
import pdfplumber
from docx import Document
from pptx import Presentation
import pickle
import logging

# Suppress pdfplumber CropBox warnings
logging.getLogger("pdfplumber").setLevel(logging.ERROR)

# Define the dataset directory
DATASET_DIR = "dataset"  # Ensure this matches your folder location


# Extraction functions
# def extract_text_from_pdf(file_path):
#     text = ""
#     try:
#         with pdfplumber.open(file_path) as pdf:
#             for page in pdf.pages:
#                 text += page.extract_text() or ""
#     except Exception as e:
#         print(f"Error extracting PDF {file_path}: {e}")
#     return text.strip()
def extract_text_from_pdf(file_path, max_pages=5):
    text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages):
                if i >= max_pages:  # Stop after `max_pages`
                    break
                text += page.extract_text() or ""
    except Exception as e:
        print(f"Error extracting PDF {file_path}: {e}")
    return text.strip()


def extract_text_from_docx(file_path):
    text = ""
    try:
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        print(f"Error extracting DOCX {file_path}: {e}")
    return text.strip()


def extract_text_from_pptx(file_path):
    text = []
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
    except Exception as e:
        print(f"Error extracting PPTX {file_path}: {e}")
    return "\n".join(text).strip()


# def extract_text(file_path):
#     if file_path.endswith(".pdf"):
#         return extract_text_from_pdf(file_path)
#     elif file_path.endswith(".docx"):
#         return extract_text_from_docx(file_path)
#     elif file_path.endswith(".pptx"):
#         return extract_text_from_pptx(file_path)
#     else:
#         print(f"Unsupported file type: {file_path}")
#         return ""
def extract_text(file_path):
    if os.path.getsize(file_path) > 5 * 1024 * 1024:  # Skip files larger than 5MB
        print(f"Skipping large file: {file_path}")
        return ""
    if file_path.endswith(".pdf"):
        return extract_text_from_pdf(file_path)
    elif file_path.endswith(".docx"):
        return extract_text_from_docx(file_path)
    elif file_path.endswith(".pptx"):
        return extract_text_from_pptx(file_path)
    else:
        print(f"Unsupported file type: {file_path}")
        return ""


# Collect texts and labels
texts = []
labels = []
DATASET_DIR = "C://Doc"

# for folder_name in os.listdir(DATASET_DIR):
#     folder_path = os.path.join(DATASET_DIR, folder_name)
#     if os.path.isdir(folder_path):
#         print(f"Processing folder: {folder_name}")
#         for file_name in os.listdir(folder_path):
#             file_path = os.path.join(folder_path, file_name)
#             if os.path.isfile(file_path):
#                 text = extract_text(file_path)
#                 if text:
#                     texts.append(text)
#                     labels.append(folder_name)
#                     print(f"Processed {file_name} -> Label: {folder_name}")
#                 else:
#                     print(f"Skipped {file_name} (no text extracted)")

import logging

logging.basicConfig(level=logging.INFO)

for folder_name in os.listdir(DATASET_DIR):
    folder_path = os.path.join(DATASET_DIR, folder_name)
    if os.path.isdir(folder_path):
        for file_name in os.listdir(folder_path):
            file_path = os.path.join(folder_path, file_name)
            if os.path.isfile(file_path):
                logging.info(f"Processing {file_path}")  # Logs the file being processed
                text = extract_text(file_path)
                if text:
                    texts.append(text)
                    labels.append(folder_name)
                    logging.info(f"Processed {file_name} -> Label: {folder_name}")
                else:
                    logging.warning(f"Skipped {file_name} (no text extracted)")


# Save the data
try:
    with open("training_data.pkl", "wb") as f:
        pickle.dump({"texts": texts, "labels": labels}, f)
    print(
        f"Collected {len(texts)} samples across {len(set(labels))} labels: {set(labels)}"
    )
except Exception as e:
    print(f"Error saving training data: {e}")
