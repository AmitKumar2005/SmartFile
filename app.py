from flask import (
    Flask,
    request,
    jsonify,
    send_from_directory,
    session,
    render_template,
    redirect,
    make_response,
)
from flask_cors import CORS
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address
from flask_wtf.csrf import CSRFProtect
import os
import logging
import pdfplumber
import fitz
from pdf2image import convert_from_path
import pytesseract
from PIL import Image
import base64
import pickle
import json
import mysql.connector
from mysql.connector import pooling
from docx import Document
from pptx import Presentation
from werkzeug.security import generate_password_hash, check_password_hash
from functools import wraps
from dotenv import load_dotenv
from datetime import timedelta
import re
import mimetypes
import secrets

# Set up logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[logging.FileHandler("app.log"), logging.StreamHandler()],
)

# Load environment variables
load_dotenv()

app = Flask(__name__)
app.secret_key = os.getenv("FLASK_SECRET_KEY")
if not app.secret_key:
    logging.error("FLASK_SECRET_KEY is not set")
    raise ValueError("FLASK_SECRET_KEY must be set in .env")
app.permanent_session_lifetime = timedelta(minutes=30)  # Session timeout
app.config["SESSION_COOKIE_SECURE"] = True
app.config["SESSION_COOKIE_HTTPONLY"] = True
app.config["SESSION_COOKIE_SAMESITE"] = "Lax"
CORS(
    app,
    resources={
        r"/*": {"origins": os.getenv("ALLOWED_ORIGINS", "*")},
    },
)
csrf = CSRFProtect(app)

# Rate limiting
limiter = Limiter(
    get_remote_address, app=app, default_limits=["200 per day", "50 per hour"]
)

DB_CONFIG = {
    "host": os.getenv("DB_HOST", "localhost"),
    "user": os.getenv("DB_USER"),
    "password": os.getenv("DB_PASSWORD"),
    "database": os.getenv("DB_NAME"),
    "pool_name": "mypool",
    "pool_size": 5,
}
for key, value in DB_CONFIG.items():
    if not value and key != "pool_name" and key != "pool_size":
        logging.error(f"Missing DB_CONFIG: {key}")
        raise ValueError(f"Missing DB_CONFIG: {key}")
MAX_TEXT_LENGTH = 1_000_000  # 1MB limit for text
ALLOWED_MIMETYPES = {
    "application/pdf": ".pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
}

# Initialize database connection pool
try:
    db_pool = pooling.MySQLConnectionPool(**DB_CONFIG)
    logging.info("MySQL connection pool initialized")
except Exception as e:
    logging.error(f"MySQL connection pool initialization failed: {e}")
    raise

# Load model and vectorizer
try:
    if not os.path.exists("model.pkl"):
        raise FileNotFoundError(
            "model.pkl not found. Run train_model.py to generate it."
        )
    if not os.path.exists("vectorizer.pkl"):
        raise FileNotFoundError(
            "vectorizer.pkl not found. Run train_model.py to generate it."
        )
    with open("model.pkl", "rb") as f:
        model = pickle.load(f)
    with open("vectorizer.pkl", "rb") as f:
        vectorizer = pickle.load(f)
    logging.info("Model and vectorizer loaded successfully")
except Exception as e:
    logging.error(f"Failed to load model or vectorizer: {e}")
    raise


def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if "user_id" not in session:
            logging.error("Unauthorized access to protected route")
            return jsonify({"error": "Unauthorized, please log in"}), 401
        return f(*args, **kwargs)

    return decorated_function


def get_email():
    """Retrieve email for logged-in user."""
    if "user_id" in session:
        try:
            conn = db_pool.get_connection()
            cursor = conn.cursor()
            cursor.execute(
                "SELECT email FROM users WHERE id = %s", (session["user_id"],)
            )
            email = cursor.fetchone()
            cursor.close()
            conn.close()
            return email[0] if email else None
        except Exception as e:
            logging.error(f"Failed to fetch email: {e}")
            return None
    return None


def validate_email(email):
    """Validate email format."""
    pattern = r"^[\w\.-]+@[\w\.-]+\.\w+$"
    return re.match(pattern, email) is not None


def sanitize_filename(filename):
    """Sanitize filename to prevent path traversal."""
    return re.sub(r"[^\w\.-]", "_", os.path.basename(filename))


def validate_file(file):
    """Validate file type and size."""
    if not file:
        return False, "No file provided"
    filename = file.filename
    if not filename:
        return False, "Empty filename"

    # Check file extension
    ext = os.path.splitext(filename)[1].lower()
    if ext not in [".pdf", ".docx", ".pptx"]:
        return False, "Invalid file extension"

    # Check MIME type
    mime_type, _ = mimetypes.guess_type(filename)
    if mime_type not in ALLOWED_MIMETYPES:
        return False, "Invalid file type"

    # Check file size
    file.seek(0, os.SEEK_END)
    file_size = file.tell()
    file.seek(0)
    if file_size > 10 * 1024 * 1024:
        return False, "File size exceeds 10MB limit"

    return True, ""


def extract_everything_from_pdf(pdf_path):
    full_text = ""
    ocr_text = ""
    extracted_images = []
    extracted_tables = []
    metadata = {}
    try:
        doc = fitz.open(pdf_path)
        if doc.is_encrypted:
            doc.close()
            logging.error("PDF is encrypted")
            return {"error": "PDF is encrypted, cannot process"}
        doc.close()
    except Exception as e:
        logging.error(f"Invalid PDF: {e}")
        return {"error": f"Invalid PDF: {str(e)}"}
    try:
        with pdfplumber.open(pdf_path) as pdf:
            metadata["page_count"] = len(pdf.pages)
            metadata["pdf_metadata"] = pdf.metadata or {}
            for i, page in enumerate(pdf.pages):
                try:
                    page_text = page.extract_text() or ""
                    if page_text.strip():
                        full_text += f"\n--- Page {i+1} ---\n{page_text}"
                    else:
                        logging.info(f"Page {i+1} is likely scanned, will use OCR")
                        ocr_text += f"\n--- OCR Page {i+1} ---\n"
                    tables = page.extract_tables()
                    for table in tables:
                        cleaned_table = [
                            [cell or "" for cell in row]
                            for row in table
                            if any(cell for cell in row)
                        ]
                        if cleaned_table:
                            extracted_tables.append(
                                {"page": i + 1, "table": cleaned_table}
                            )
                except Exception as e:
                    logging.error(f"Error processing page {i+1} with pdfplumber: {e}")
    except Exception as e:
        logging.error(f"Error with pdfplumber: {e}")
        return {"error": f"pdfplumber error: {str(e)}"}
    try:
        images = convert_from_path(pdf_path, dpi=200)
        for i, img in enumerate(images):
            if (
                not full_text.split(f"--- Page {i+1} ---")[1].strip()
                if f"--- Page {i+1} ---" in full_text
                else True
            ):
                try:
                    img = img.convert("L")
                    ocr_result = pytesseract.image_to_string(img, config="--psm 6")
                    ocr_text += ocr_result.strip() + "\n"
                except Exception as e:
                    logging.error(f"OCR error on page {i+1}: {e}")
                finally:
                    img.close()
    except Exception as e:
        logging.error(f"Error converting PDF to images: {e}")
    try:
        doc = fitz.open(pdf_path)
        for page_index in range(len(doc)):
            image_list = doc.get_page_images(page_index)
            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    ext = base_image["ext"]
                    image_b64 = base64.b64encode(image_bytes).decode("utf-8")
                    extracted_images.append(
                        {
                            "page": page_index + 1,
                            "img_index": img_index + 1,
                            "ext": ext,
                            "data": image_b64,
                        }
                    )
                except Exception as e:
                    logging.error(
                        f"Error extracting image {img_index+1} on page {page_index+1}: {e}"
                    )
        doc.close()
    except Exception as e:
        logging.error(f"Error with PyMuPDF: {e}")
    full_text = "\n".join(
        line.strip() for line in full_text.splitlines() if line.strip()
    )
    ocr_text = "\n".join(line.strip() for line in ocr_text.splitlines() if line.strip())
    table_text = "\n".join(
        "\t".join(cell for cell in row)
        for table in extracted_tables
        for row in table["table"]
    )
    combined_text = f"{full_text}\n{ocr_text}\n{table_text}".strip()
    return {
        "text": full_text,
        "ocr_text": ocr_text,
        "combined_text": combined_text,
        "images": extracted_images,
        "tables": extracted_tables,
        "metadata": metadata,
    }


def extract_text_from_docx(file_path):
    text = ""
    try:
        doc = Document(file_path)
        text = "\n".join(para.text for para in doc.paragraphs if para.text.strip())
    except Exception as e:
        logging.error(f"Failed to read DOCX {file_path}: {e}")
        return {"error": f"DOCX error: {str(e)}"}
    return {
        "text": text,
        "combined_text": text,
        "images": [],
        "tables": [],
        "metadata": {"file_type": "docx"},
    }


def extract_text_from_pptx(file_path):
    text = []
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
    except Exception as e:
        logging.error(f"Failed to read PPTX {file_path}: {e}")
        return {"error": f"PPTX error: {str(e)}"}
    combined_text = "\n".join(text).strip()
    return {
        "text": combined_text,
        "combined_text": combined_text,
        "images": [],
        "tables": [],
        "metadata": {"file_type": "pptx", "slide_count": len(prs.slides)},
    }


def extract_content(file_path):
    if file_path.lower().endswith(".pdf"):
        return extract_everything_from_pdf(file_path)
    elif file_path.lower().endswith(".docx"):
        return extract_text_from_docx(file_path)
    elif file_path.lower().endswith(".pptx"):
        return extract_text_from_pptx(file_path)
    else:
        return {"error": "Unsupported file type"}


@app.route("/favicon.ico")
def favicon():
    return send_from_directory(os.path.join(app.root_path, "static"), "favicon.ico")


@app.route("/static/favicon.ico")
def static_favicon():
    return send_from_directory(os.path.join(app.root_path, "static"), "favicon.ico")


@app.route("/")
def index():
    logged_in = "user_id" in session
    username = get_email() if logged_in else None
    logging.info(
        f"Index route - Session: {session}, Logged in: {logged_in}, Username: {username}"
    )
    response = make_response(
        render_template("index.html", logged_in=logged_in, username=username)
    )
    response.headers["Cache-Control"] = "no-store, no-cache, must-revalidate, max-age=0"
    return response


@app.route("/login")
def login_page():
    logged_in = "user_id" in session
    username = get_email() if logged_in else None
    if logged_in:
        return redirect("/upload")
    response = make_response(
        render_template("login.html", logged_in=logged_in, username=username)
    )
    response.headers["Cache-Control"] = "no-store, no-cache, must-revalidate, max-age=0"
    return response


@app.route("/register")
def register_page():
    logged_in = "user_id" in session
    username = get_email() if logged_in else None
    if logged_in:
        return redirect("/upload")
    response = make_response(
        render_template("register.html", logged_in=logged_in, username=username)
    )
    response.headers["Cache-Control"] = "no-store, no-cache, must-revalidate, max-age=0"
    return response


@app.route("/upload")
@login_required
def upload_page():
    logged_in = "user_id" in session
    username = get_email()
    response = make_response(
        render_template("upload.html", logged_in=logged_in, username=username)
    )
    response.headers["Cache-Control"] = "no-store, no-cache, must-revalidate, max-age=0"
    return response


@app.route("/search")
@login_required
def search_page():
    logged_in = "user_id" in session
    username = get_email()
    response = make_response(
        render_template("search.html", logged_in=logged_in, username=username)
    )
    response.headers["Cache-Control"] = "no-store, no-cache, must-revalidate, max-age=0"
    return response


@app.route("/debug_session")
def debug_session():
    return jsonify(dict(session))


@app.route("/register", methods=["POST"])
@limiter.limit("10 per minute")
@csrf.exempt
def register():
    data = request.form
    if not data or "email" not in data or "password" not in data:
        logging.error("Missing email or password")
        return jsonify({"error": "Email and password required"}), 400
    email = data["email"].strip()
    password = data["password"]
    if not validate_email(email):
        logging.error("Invalid email format")
        return jsonify({"error": "Invalid email format"}), 400
    if len(password) < 8:
        logging.error("Password too short")
        return jsonify({"error": "Password must be at least 8 characters"}), 400
    try:
        conn = db_pool.get_connection()
        cursor = conn.cursor()
        cursor.execute("SELECT id FROM users WHERE email = %s", (email,))
        if cursor.fetchone():
            cursor.close()
            conn.close()
            logging.error(f"Email {email} already registered")
            return jsonify({"error": "This email is already registered"}), 400
        hashed_password = generate_password_hash(password)
        cursor.execute(
            "INSERT INTO users (email, password) VALUES (%s, %s)",
            (email, hashed_password),
        )
        conn.commit()
        cursor.execute("SELECT id FROM users WHERE email = %s", (email,))
        user_id = cursor.fetchone()[0]
        session["user_id"] = user_id
        session.permanent = True
        session["csrf_token"] = secrets.token_hex(16)
        cursor.close()
        conn.close()
        logging.info(
            f"User {email} registered and logged in successfully, Session: {session}"
        )
        response = make_response(jsonify({"message": "Registration successful"}))
        response.headers["HX-Redirect"] = "/"
        response.headers["Cache-Control"] = (
            "no-store, no-cache, must-revalidate, max-age=0"
        )
        return response
    except Exception as e:
        logging.error(f"Registration error: {e}")
        return jsonify({"error": "Registration failed"}), 500


@app.route("/login", methods=["POST"])
@limiter.limit("10 per minute")
@csrf.exempt
def login():
    data = request.form
    logging.info(f"Login attempt with data: {data}")
    if not data or "email" not in data or "password" not in data:
        logging.error("Missing email or password")
        return jsonify({"error": "Email and password required"}), 400
    email = data["email"].strip()
    password = data["password"]
    try:
        conn = db_pool.get_connection()
        cursor = conn.cursor()
        cursor.execute("SELECT id, password FROM users WHERE email = %s", (email,))
        user = cursor.fetchone()
        cursor.close()
        conn.close()
        if user and check_password_hash(user[1], password):
            session["user_id"] = user[0]
            session.permanent = True
            session["csrf_token"] = secrets.token_hex(16)
            logging.info(f"Login successful for user_id: {user[0]}, Session: {session}")
            response = make_response(jsonify({"message": "Login successful"}))
            response.headers["HX-Redirect"] = "/"
            response.headers["Cache-Control"] = (
                "no-store, no-cache, must-revalidate, max-age=0"
            )
            return response
        logging.error(f"Invalid credentials for email: {email}")
        return jsonify({"error": "Invalid Email or Password"}), 401
    except Exception as e:
        logging.error(f"Login error: {e}")
        return jsonify({"error": "Login failed, please try again"}), 500


@app.route("/logout", methods=["POST"])
@csrf.exempt
def logout():
    session.pop("user_id", None)
    session.pop("csrf_token", None)
    logging.info("User logged out")
    response = make_response(jsonify({"message": "Logged out successfully"}))
    response.set_cookie(
        "session",
        "",
        expires=0,
        secure=app.config["SESSION_COOKIE_SECURE"],
        httponly=app.config["SESSION_COOKIE_HTTPONLY"],
        samesite=app.config["SESSION_COOKIE_SAMESITE"],
    )
    response.headers["HX-Redirect"] = "/"
    response.headers["Cache-Control"] = "no-store, no-cache, must-revalidate, max-age=0"
    return response


# @app.route("/extract", methods=["POST"])
# @login_required
# @limiter.limit("5 per minute")
# @csrf.exempt
# def extract():
#     if "file" not in request.files:
#         logging.error("No file uploaded")
#         return jsonify({"error": "No file uploaded"}), 400
#     file = request.files["file"]
#     is_valid, error_message = validate_file(file)
#     if not is_valid:
#         logging.error(error_message)
#         return jsonify({"error": error_message}), 400

#     user_id = session["user_id"]
#     sanitized_filename = sanitize_filename(file.filename)

#     # Check if file already exists for this user
#     try:
#         conn = db_pool.get_connection()
#         cursor = conn.cursor()
#         cursor.execute(
#             "SELECT id FROM files WHERE user_id = %s AND filename = %s",
#             (user_id, sanitized_filename),
#         )
#         if cursor.fetchone():
#             cursor.close()
#             conn.close()
#             logging.info(f"File {sanitized_filename} already exists for user {user_id}")
#             return jsonify({"error": "File already available"}), 200
#         cursor.close()
#         conn.close()
#     except Exception as e:
#         logging.error(f"Error checking for existing file: {e}")
#         return jsonify({"error": "Failed to check file existence"}), 500

#     user_dir = os.path.join("Uploads", f"user_{user_id}")
#     os.makedirs(user_dir, exist_ok=True)
#     temp_path = os.path.join(user_dir, sanitized_filename)
#     file.save(temp_path)
#     logging.info(f"File saved temporarily: {temp_path}")

#     result = extract_content(temp_path)
#     if "error" in result:
#         logging.error(f"Extraction failed: {result['error']}")
#         os.remove(temp_path)
#         return jsonify({"error": result["error"]}), 400

#     combined_text = result.get("combined_text", "")
#     predicted_folder = "Unknown"
#     if combined_text.strip():
#         try:
#             X = vectorizer.transform([combined_text])
#             predicted_folder = model.predict(X)[0]
#             result["predicted_folder"] = predicted_folder
#         except Exception as e:
#             logging.error(f"Prediction failed: {e}")

#     folder_dir = os.path.join(user_dir, predicted_folder)
#     os.makedirs(folder_dir, exist_ok=True)
#     final_path = os.path.join(folder_dir, sanitized_filename)
#     os.rename(temp_path, final_path)
#     logging.info(f"File moved to: {final_path}")

#     if len(combined_text.encode("utf-8")) > MAX_TEXT_LENGTH:
#         logging.warning(
#             f"Text for {sanitized_filename} exceeds {MAX_TEXT_LENGTH} bytes, truncating."
#         )
#         combined_text = combined_text.encode("utf-8")[:MAX_TEXT_LENGTH].decode(
#             "utf-8", errors="ignore"
#         )

#     try:
#         conn = db_pool.get_connection()
#         cursor = conn.cursor()
#         cursor.execute(
#             "INSERT INTO files (user_id, filename, folder, text, tables, file_path) VALUES (%s, %s, %s, %s, %s, %s)",
#             (
#                 user_id,
#                 sanitized_filename,
#                 predicted_folder,
#                 combined_text,
#                 json.dumps(result.get("tables", [])),
#                 final_path,
#             ),
#         )
#         file_id = cursor.lastrowid
#         conn.commit()
#         cursor.close()
#         conn.close()
#     except mysql.connector.Error as e:
#         logging.error(f"Failed to save file {sanitized_filename} to database: {e}")
#         os.remove(final_path)
#         return jsonify({"error": "Failed to save file metadata"}), 500

#     logging.info(f"Extraction successful. Predicted folder: {predicted_folder}")
#     return jsonify(
#         {
#             "file_id": file_id,
#             "filename": sanitized_filename,
#             "predicted_folder": predicted_folder,
#             "text": combined_text,
#             "tables": result.get("tables", []),
#         }
#     )


@app.route("/extract", methods=["POST"])
@login_required
@limiter.limit("5 per minute")
@csrf.exempt
def extract():
    if "file" not in request.files:
        logging.error("No file uploaded")
        return jsonify({"error": "No file uploaded"}), 400
    file = request.files["file"]
    is_valid, error_message = validate_file(file)
    if not is_valid:
        logging.error(error_message)
        return jsonify({"error": error_message}), 400

    user_id = session["user_id"]
    sanitized_filename = sanitize_filename(file.filename)

    # Check if file already exists for this user
    try:
        conn = db_pool.get_connection()
        cursor = conn.cursor()
        cursor.execute(
            "SELECT id FROM files WHERE user_id = %s AND filename = %s",
            (user_id, sanitized_filename),
        )
        if cursor.fetchone():
            cursor.close()
            conn.close()
            logging.info(f"File {sanitized_filename} already exists for user {user_id}")
            return jsonify({"error": "File already available"}), 200
        cursor.close()
        conn.close()
    except Exception as e:
        logging.error(f"Error checking for existing file: {e}")
        return jsonify({"error": "Failed to check file existence"}), 500

    user_dir = os.path.join("Uploads", f"user_{user_id}")
    os.makedirs(user_dir, exist_ok=True)
    temp_path = os.path.join(user_dir, sanitized_filename)
    file.save(temp_path)
    logging.info(f"File saved temporarily: {temp_path}")

    result = extract_content(temp_path)
    if "error" in result:
        logging.error(f"Extraction failed: {result['error']}")
        os.remove(temp_path)
        return jsonify({"error": result["error"]}), 400

    combined_text = result.get("combined_text", "")
    predicted_folder = "Unknown"
    if combined_text.strip():
        try:
            X = vectorizer.transform([combined_text])
            predicted_folder = model.predict(X)[0]
            result["predicted_folder"] = predicted_folder
        except Exception as e:
            logging.error(f"Prediction failed: {e}")

    folder_dir = os.path.join(user_dir, predicted_folder)
    os.makedirs(folder_dir, exist_ok=True)
    final_path = os.path.join(folder_dir, sanitized_filename)
    logging.info(f"Moving file from {temp_path} to {final_path}")
    try:
        os.rename(temp_path, final_path)
        if not os.path.exists(final_path):
            raise FileNotFoundError(f"File not found at {final_path} after move")
        logging.info(f"File successfully moved to {final_path}")
    except Exception as e:
        logging.error(f"Failed to move file to {final_path}: {e}")
        os.remove(temp_path)  # Clean up temp file
        return jsonify({"error": f"Failed to move file: {str(e)}"}), 500

    if len(combined_text.encode("utf-8")) > MAX_TEXT_LENGTH:
        logging.warning(
            f"Text for {sanitized_filename} exceeds {MAX_TEXT_LENGTH} bytes, truncating."
        )
        combined_text = combined_text.encode("utf-8")[:MAX_TEXT_LENGTH].decode(
            "utf-8", errors="ignore"
        )

    try:
        conn = db_pool.get_connection()
        cursor = conn.cursor()
        cursor.execute(
            "INSERT INTO files (user_id, filename, folder, text, tables, file_path) VALUES (%s, %s, %s, %s, %s, %s)",
            (
                user_id,
                sanitized_filename,
                predicted_folder,
                combined_text,
                json.dumps(result.get("tables", [])),
                final_path,
            ),
        )
        file_id = cursor.lastrowid
        conn.commit()
        cursor.close()
        conn.close()
    except mysql.connector.Error as e:
        logging.error(f"Failed to save file {sanitized_filename} to database: {e}")
        os.remove(final_path)
        return jsonify({"error": "Failed to save file metadata"}), 500

    logging.info(f"Extraction successful. Predicted folder: {predicted_folder}")
    return jsonify(
        {
            "file_id": file_id,
            "filename": sanitized_filename,
            "predicted_folder": predicted_folder,
            "text": combined_text,
            "tables": result.get("tables", []),
        }
    )


@app.route("/search", methods=["POST"])
@login_required
@limiter.limit("20 per minute")
@csrf.exempt
def search():
    data = request.form
    if not data or "query" not in data:
        logging.error("No query provided")
        return "<div class='error-message'>No query provided</div>", 400
    query = data["query"].strip()[:255]  # Limit query length
    if not query:
        logging.error("Empty query")
        return "<div class='error-message'>Empty query</div>", 400
    user_id = session["user_id"]
    try:
        conn = db_pool.get_connection()
        cursor = conn.cursor()
        cursor.execute(
            "SELECT id, filename, folder, text FROM files WHERE user_id = %s AND MATCH(text) AGAINST(%s IN BOOLEAN MODE) LIMIT 20",
            (user_id, query),
        )
        results = [
            {
                "id": row[0],
                "filename": row[1],
                "folder": row[2],
                "snippet": row[3][:200] + "..." if row[3] else "",
            }
            for row in cursor.fetchall()
        ]
        cursor.close()
        conn.close()
        logging.info(
            f"Search completed. Found {len(results)} matches for query: {query}"
        )
        if not results:
            return "<div class='noFolder'>No results found</div>", 200
        html = "<div class='search-result-grid'>"
        for result in results:
            html += f"""
                <div class='search-result-card'>
                    <a href='/view_file/{result['id']}' target='_blank' class='link font-semibold'>{result['filename']}</a>
                    <p>Folder: {result['folder']}</p>
                </div>
            """
        html += "</div>"
        return html, 200
    except Exception as e:
        logging.error(f"Search error: {e}")
        return (
            f"<div class='error-message'>Search failed: {str(e)}</div>",
            500,
        )


@app.route("/list_files", methods=["GET"])
@login_required
def list_files():
    user_id = session["user_id"]
    try:
        conn = db_pool.get_connection()
        cursor = conn.cursor()
        cursor.execute(
            "SELECT DISTINCT folder FROM files WHERE user_id = %s ORDER BY folder",
            (user_id,),
        )
        folders = [row[0] for row in cursor.fetchall()]
        cursor.close()
        conn.close()
        html = "<div class='file-list-grid'>"
        for folder in folders:
            html += f"""
                <div class='folder' 
                     hx-get='/folder_files/{folder}' 
                     hx-target='#file-list' 
                     hx-swap='innerHTML'>
                    <h3 class='header'>{folder}</h3>
                </div>
            """
        html += "</div>"
        if not folders:
            html = "<div class='noFolder'>No folders or files uploaded yet.</div>"
        return html, 200
    except Exception as e:
        logging.error(f"List files error: {e}")
        return (
            "<div class='failToLoad'>Failed to load folders</div>",
            500,
        )


@app.route("/folder_files/<folder>", methods=["GET"])
@login_required
def folder_files(folder):
    user_id = session["user_id"]
    # folder = re.sub(r"[^\w\-]", "", folder)  # Sanitize folder name
    logging.info(folder)
    try:
        conn = db_pool.get_connection()
        cursor = conn.cursor()
        cursor.execute(
            "SELECT id, filename, folder FROM files WHERE user_id = %s AND folder = %s ORDER BY filename",
            (user_id, folder),
        )
        files = [
            {"id": row[0], "filename": row[1], "folder": row[2]}
            for row in cursor.fetchall()
        ]
        cursor.close()
        conn.close()
        html = f"""
            <div class='card'>
                <h3 class='text-lg'>Files in {folder}</h3>
                <button hx-get='/list_files' hx-target='#file-list' hx-swap='innerHTML' 
                        class='link mb-6'>Back to Folders</button>
                <ul class='list'>
        """
        for file in files:
            html += f"""
                <li class='mb-4'>
                    <a href='/view_file/{file["id"]}' target='_blank' class='link'>{file["filename"]}</a>
                </li>
            """
        html += "</ul></div>"
        if not files:
            html = f"""
                <div class='card'>
                    <h3 class='text-lg'>Files in {folder}</h3>
                    <button hx-get='/list_files' hx-target='#file-list' hx-swap='innerHTML' 
                            class='link mb-6'>Back to Folders</button>
                    <div class='noFolder'>No files in this folder.</div>
                </div>
            """
        return html, 200
    except Exception as e:
        logging.error(f"Folder files error: {e}")
        return "<div class='failToLoad'>Failed to load files</div>", 500


@app.route("/view_file/<int:file_id>", methods=["GET"])
@login_required
def view_file(file_id):
    user_id = session["user_id"]
    try:
        conn = db_pool.get_connection()
        cursor = conn.cursor()
        cursor.execute(
            "SELECT filename, file_path, text, tables FROM files WHERE id = %s AND user_id = %s",
            (file_id, user_id),
        )
        file = cursor.fetchone()
        cursor.close()
        conn.close()
        if not file:
            return jsonify({"error": "File not found or access denied"}), 404
        filename, file_path, text, tables = file
        if os.path.exists(file_path):
            return send_from_directory(
                os.path.dirname(file_path), os.path.basename(file_path)
            )
        tables_data = json.loads(tables) if tables else []
        html = f"""
            <html>
            <head>
                <title>{filename}</title>
                <link href="/static/styles.css" rel="stylesheet">
            </head>
            <body>
                <h1>{filename}</h1>
                <h2>Extracted Text</h2>
                <pre>{text}</pre>
                <h2>Tables</h2>
        """
        for table in tables_data:
            html += "<table>"
            for row in table["table"]:
                html += "<tr>"
                for cell in row:
                    html += f"<td>{cell}</td>"
                html += "</tr>"
            html += "</table>"
        html += "</body></html>"
        return html, 200
    except Exception as e:
        logging.error(f"View file error: {e}")
        return jsonify({"error": "Failed to view file"}), 500


@app.route("/feedback", methods=["POST"])
@login_required
@csrf.exempt
def feedback():
    data = request.form
    file_id = data.get("file_id")
    corrected_folder = data.get("corrected_folder")
    user_id = session["user_id"]
    try:
        conn = db_pool.get_connection()
        cursor = conn.cursor()
        cursor.execute(
            "SELECT folder, file_path FROM files WHERE id = %s AND user_id = %s",
            (file_id, user_id),
        )
        result = cursor.fetchone()
        if not result:
            cursor.close()
            conn.close()
            return jsonify({"error": "File not found"}), 404
        predicted_folder, file_path = result
        cursor.execute(
            "INSERT INTO feedback (file_id, user_id, predicted_folder, corrected_folder) VALUES (%s, %s, %s, %s)",
            (file_id, user_id, predicted_folder, corrected_folder),
        )
        if corrected_folder and corrected_folder != predicted_folder:
            cursor.execute(
                "UPDATE files SET folder = %s WHERE id = %s",
                (corrected_folder, file_id),
            )
            new_folder = os.path.join(
                os.path.dirname(os.path.dirname(file_path)), corrected_folder
            )
            os.makedirs(new_folder, exist_ok=True)
            new_path = os.path.join(new_folder, os.path.basename(file_path))
            os.rename(file_path, new_path)
            cursor.execute(
                "UPDATE files SET file_path = %s WHERE id = %s", (new_path, file_id)
            )
        conn.commit()
        cursor.close()
        conn.close()
        logging.info(f"Feedback recorded for file_id: {file_id}")
        return "<div class='success-message'>Feedback submitted</div>", 200
    except Exception as e:
        logging.error(f"Feedback error: {e}")
        return (
            "<div class='error-message'>Failed to submit feedback</div>",
            500,
        )


if __name__ == "__main__":
    logging.info("Starting Flask application")
    app.run(debug=True, port=5000, host="0.0.0.0")
