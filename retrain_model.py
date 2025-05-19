import mysql.connector
import pickle
import logging
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.linear_model import LogisticRegression
from dotenv import load_dotenv
import os
import shutil

load_dotenv()

logging.basicConfig(
    level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s"
)

DB_CONFIG = {
    "host": os.getenv("DB_HOST"),
    "port": os.getenv("DB_PORT", "4000"),
    "user": os.getenv("DB_USER"),
    "password": os.getenv("DB_PASSWORD"),
    "database": os.getenv("DB_NAME"),
    "ssl_ca": os.path.join(os.path.dirname(__file__), "tidb-ca.pem"),
    "ssl_verify_cert": True,
    "ssl_verify_identity": True,
    "use_pure": True,
}

TEMP_DIR = "/tmp/TempTraining"


def extract_content_from_pdf(file_path):
    import pdfplumber

    text = ""
    tables = []
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text() or ""
                text += page_text + "\n"
                page_tables = page.extract_tables()
                for table in page_tables:
                    cleaned_table = [
                        [cell or "" for cell in row]
                        for row in table
                        if any(cell for cell in row)
                    ]
                    if cleaned_table:
                        tables.append(
                            {"page": page.page_number, "table": cleaned_table}
                        )
    except Exception as e:
        logging.error(f"Failed to read PDF {file_path}: {e}")
    return text.strip(), tables


def extract_text_from_docx(file_path):
    from docx import Document

    text = ""
    try:
        doc = Document(file_path)
        text = "\n".join(para.text for para in doc.paragraphs if para.text.strip())
    except Exception as e:
        logging.error(f"Failed to read DOCX {file_path}: {e}")
    return text.strip(), []


def extract_text_from_pptx(file_path):
    from pptx import Presentation

    text = []
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
    except Exception as e:
        logging.error(f"Failed to read PPTX {file_path}: {e}")
    return "\n".join(text).strip(), []


def extract_content(file_path):
    if file_path.lower().endswith(".pdf"):
        return extract_content_from_pdf(file_path)
    elif file_path.lower().endswith(".docx"):
        return extract_text_from_docx(file_path)
    elif file_path.lower().endswith(".pptx"):
        return extract_text_from_pptx(file_path)
    else:
        logging.warning(f"Unsupported file type: {file_path}")
        return "", []


def retrain_model():
    try:
        conn = mysql.connector.connect(**DB_CONFIG)
        cursor = conn.cursor()
        cursor.execute(
            """
            SELECT f.file_id, f.corrected_folder, fi.file_path
            FROM feedback f
            JOIN files fi ON f.file_id = fi.id
            WHERE f.corrected_folder IS NOT NULL
        """
        )
        feedback_data = cursor.fetchall()
        cursor.close()
        conn.close()

        if not feedback_data:
            logging.info("No feedback data to retrain model")
            return

        if os.path.exists(TEMP_DIR):
            shutil.rmtree(TEMP_DIR)
        os.makedirs(TEMP_DIR)

        texts = []
        labels = []
        for file_id, corrected_folder, file_path in feedback_data:
            temp_folder = os.path.join(TEMP_DIR, corrected_folder)
            os.makedirs(temp_folder, exist_ok=True)
            temp_path = os.path.join(temp_folder, os.path.basename(file_path))
            shutil.copy2(file_path, temp_path)
            text, _ = extract_content(temp_path)
            if text:
                texts.append(text)
                labels.append(corrected_folder)

        shutil.rmtree(TEMP_DIR)

        if not texts:
            logging.info("No valid feedback data to retrain model")
            return

        vectorizer = TfidfVectorizer(max_features=5000, stop_words="english")
        X = vectorizer.fit_transform(texts)
        model = LogisticRegression(max_iter=1000, C=1.0)
        model.fit(X, labels)
        with open("model.pkl", "wb") as f:
            pickle.dump(model, f)
        with open("vectorizer.pkl", "wb") as f:
            pickle.dump(vectorizer, f)
        logging.info("Model retrained and saved")
    except Exception as e:
        logging.error(f"Retrain model error: {e}")


if __name__ == "__main__":
    retrain_model()
